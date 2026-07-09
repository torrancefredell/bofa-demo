#!/usr/bin/env python3
"""Build BofA x Cognition demo deck (16:9)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy
import os

NAVY = RGBColor(0x01, 0x21, 0x69)   # BofA blue
RED = RGBColor(0xE3, 0x18, 0x37)    # BofA red
LIGHT = RGBColor(0xF4, 0xF6, 0xF9)
GRAY = RGBColor(0x5A, 0x64, 0x72)
DARK = RGBColor(0x1C, 0x24, 0x30)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MIDBLUE = RGBColor(0x2E, 0x5B, 0xB8)
FONT = "Roboto"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = 13.333, 7.5

def add_slide():
    return prs.slides.add_slide(BLANK)

def rect(slide, x, y, w, h, fill, line=None, shape=MSO_SHAPE.RECTANGLE, radius=None):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp

def tb(slide, x, y, w, h, runs, size=14, color=DARK, bold=False, align=PP_ALIGN.LEFT,
       anchor=MSO_ANCHOR.TOP, space_after=6, line_spacing=1.0):
    """runs: str, or list of paragraphs; each paragraph is str or (text, dict) or list of runs."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    if isinstance(runs, str):
        runs = [runs]
    first = True
    for para in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        if isinstance(para, tuple):
            text, opts = para
            rlist = [(text, opts)]
        elif isinstance(para, list):
            rlist = para
        else:
            rlist = [(para, {})]
        for text, opts in rlist:
            r = p.add_run(); r.text = text
            f = r.font
            f.name = FONT
            f.size = Pt(opts.get("size", size))
            f.bold = opts.get("bold", bold)
            f.italic = opts.get("italic", False)
            f.color.rgb = opts.get("color", color)
        if isinstance(para, tuple) and para[1].get("align"):
            p.alignment = para[1]["align"]
    return box

def header(slide, kicker, title, num):
    rect(slide, 0, 0, SW, 1.28, NAVY)
    rect(slide, 0, 1.28, SW, 0.055, RED)
    tb(slide, 0.55, 0.16, 10.5, 0.35, kicker.upper(), size=11, color=RGBColor(0x9F,0xB4,0xE0), bold=True)
    tsize = 27 if len(title) <= 44 else 22
    tb(slide, 0.55, 0.42, 11.4, 0.8, title, size=tsize, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    tb(slide, 12.3, 0.42, 0.7, 0.8, num, size=14, color=RGBColor(0x9F,0xB4,0xE0), align=PP_ALIGN.RIGHT)

def footer(slide):
    tb(slide, 0.55, 7.08, 8, 0.32, "Bank of America  ×  Cognition", size=9, color=GRAY)

def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def card(slide, x, y, w, h, title, body_lines, title_color=NAVY, accent=RED, body_size=12):
    rect(slide, x, y, w, h, LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06)
    rect(slide, x, y + 0.12, 0.06, h - 0.24, accent)
    tb(slide, x + 0.22, y + 0.12, w - 0.4, 0.4, title, size=14, color=title_color, bold=True)
    paras = []
    for ln in body_lines:
        paras.append(ln)
    tb(slide, x + 0.22, y + 0.52, w - 0.4, h - 0.65, paras, size=body_size, color=DARK, space_after=5, line_spacing=1.05)

# ---------------- Slide 1: Title ----------------
s = add_slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, 4.62, SW, 0.06, RED)
tb(s, 0.9, 2.55, 11.6, 1.2, [
    ("Bank of America  ×  Cognition", {"size": 44, "bold": True, "color": WHITE}),
])
tb(s, 0.9, 3.6, 11.0, 0.6, "A working session with Devin, the AI software engineer.", size=17, color=RGBColor(0xC9,0xD4,0xEC))
tb(s, 0.9, 4.95, 11.0, 1.2, [
    ("[Your Name]  ·  Forward Deployed Engineer, Cognition", {"size": 15, "color": WHITE, "bold": True}),
    ("[Date]", {"size": 13, "color": RGBColor(0xC9,0xD4,0xEC)}),
])
notes(s, """Keep this warm and casual — thank them for the time before you touch the deck.

"Thanks for making time — I know 45 minutes on your calendars is expensive, so I promise to spend most of it on your actual problem, not on marketing slides."

Personal hook (say it here, expand on slide 3): "Quick context on me — before Cognition I was a software engineer at Capital One, so I've lived inside a big bank's SDLC. A lot of what we'll talk about today is stuff I personally lost weekends to."

[RESEARCH BEFORE MEETING]
- Confirm exact Angular 14 EOL date your security policy keys off (community EOL was Nov 2023; confirm whether BofA has extended support via HeroDevs/OpenLogic or an internal exception process — this changes the urgency framing).
- Fill in your name/date, and get the attendees' actual names ahead of time — greet by name.

[PAUSE] after intro — ask: "Before I dive in: has anything changed since you reached out after the conference? Anything you want to make sure we cover?" Write down what they say and visibly come back to it later — that's how you keep opinionated execs engaged.""")

# ---------------- Slide 2: Agenda ----------------
s = add_slide()
header(s, "So you know where we're going", "How we'll spend the next 45 minutes", "2")
rows = [
    ("1", "Why I'm here", "~2 min", "Brief context, then straight to your problem"),
    ("2", "Your problem — and why BofA", "~7 min", "Legacy debt, compliance clocks — you talk, I listen"),
    ("3", "What Devin is", "~5 min", "The honest 2-minute version"),
    ("4", "Live demo", "~15 min", "A real framework migration, live"),
    ("5", "What it means for each of you", "~10 min", "Capacity · audit · architecture"),
    ("6", "Next steps", "~5 min", "A concrete pilot proposal"),
]
y = 1.62
for n, t, tm, d in rows:
    rect(s, 0.55, y, 12.25, 0.78, LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.12)
    rect(s, 0.72, y + 0.16, 0.46, 0.46, NAVY, shape=MSO_SHAPE.OVAL)
    tb(s, 0.72, y + 0.2, 0.46, 0.4, n, size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    tb(s, 1.38, y + 0.09, 4.6, 0.4, t, size=15, color=NAVY, bold=True)
    tb(s, 1.38, y + 0.42, 9.6, 0.32, d, size=10.5, color=GRAY)
    tb(s, 11.35, y + 0.2, 1.3, 0.4, tm, size=12, color=RED, bold=True, align=PP_ALIGN.RIGHT)
    y += 0.88
footer(s)
notes(s, """"Here's the flight plan. Two things to call out: the demo is real — live product, on a codebase structured like yours, not a video. And a third of this meeting is reserved for you talking, not me."

[ASK] "Does this split work for everyone? If you'd rather go straight to the demo, I'm happy to reshuffle."

DERAIL-PREP: This slide is your steering wheel. If the conversation goes sideways later, point back here: "Great question — that's exactly section 5, can I park it for 10 minutes so I can answer it properly?" Parking with a named return time keeps control without shutting anyone down.

--- "WHY I'M HERE" — delivered verbally off this slide, ~2 min, no dedicated slide. Keep it brief; you represent Cognition, not your biography:

"Quick context on me before we dig in: before Cognition I was a software engineer at Capital One, so I've been on your side of the table — I was the engineer who got voluntold for framework migrations. Spreadsheets of deprecations, chasing a dozen teams for sign-off on shared library changes, compliance deadlines eating weekends. So when I show you Devin today, it's not abstract for me — this is the work I wished someone would take off my plate."

Don't make the full why-BofA case here — that now has its own slide (slide 4, right after the problem slide, part of agenda item 2). Just tee it up in one breath: "And I'll say a word in a minute about why I think BofA specifically is the right place for this."

[ASK the room] "Who owns the Angular upgrade today — and what did they stop doing to own it?" Then actually wait. The VP of Engineering will usually answer. If they name a person/team, refer back to them later ("so Priya's team gets their roadmap back").

[RESEARCH BEFORE MEETING]
- Verify the current BofA tech budget figure (~$12–13B cited publicly; ~$3.5–4B on new initiatives) from their most recent earnings call.
- Erica stats (2B+ interactions cited publicly), AI patent counts, exec quotes (Aditya Bhasin, Hari Gopalkrishnan) on GenAI for developers.
- Keep the Capital One story factual — they may probe.""")

# ---------------- Slide 3: The problem ----------------
s = add_slide()
header(s, "Let's get straight to it", "The problem: legacy debt is eating your roadmap", "3")
tb(s, 0.55, 1.5, 12.25, 0.55, [[("Frameworks age out ", {"bold": True, "color": RED, "size": 17}), ("— and running unsupported software in production violates your own security policy.", {"size": 17})]])
c_w, c_h = 3.95, 2.35
card(s, 0.55, 2.3, c_w, c_h, "The debt", [
    "Outdated frameworks across the estate",
    "Every deferred upgrade compounds",
], body_size=14)
card(s, 4.7, 2.3, c_w, c_h, "The clock", [
    "Compliance-driven EOL deadlines",
    "Audit exposure while you wait",
], body_size=14)
card(s, 8.85, 2.3, c_w, c_h, "The cost", [
    "Backlog grows, greenfield stalls",
    "Best engineers stuck on upkeep",
], body_size=14)
rect(s, 0.55, 4.95, 12.25, 0.95, LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1)
tb(s, 0.85, 5.18, 11.7, 0.65, [[("Every quarter spent on maintenance is a quarter not spent building.", {"bold": True, "color": NAVY})]], size=16)
footer(s)
notes(s, """This is the most important non-demo slide. Keep it general — the specific Angular 14→18 story is the demo's job; here you're naming the pattern they live with every day.

Talk track: "Here's the pattern I suspect you know well. Frameworks and platforms age out — Angular, Java versions, Spring, on-prem stacks — and your own security policy says unsupported software can't run in production. So upgrades arrive with hard compliance deadlines, they land on your most experienced engineers, and every one of them displaces greenfield work. The backlog doesn't shrink; it compounds. The cruel part: the better your engineers are, the more of them you burn on maintenance nobody will ever see."

Bridge to the demo (one line, don't dwell): "Today's concrete example is your Angular 14→18 migration — hard EOL deadline, millions of customers, a shared library with 12+ downstream teams. We'll dig into that live in the demo."

Spend real time here and let THEM correct you — being corrected is a win, it means engagement and better intel.

[ASK, by persona]
- VP Eng: "How much of your team's capacity goes to maintenance and upgrades versus new build? What's the number you wish it was?"
- Security: "What does your policy actually require when something goes EOL — remediation plan? exception filing? What's your audit exposure today?"
- Architect: "Where's the next EOL clock ticking in your estate — and how does that upgrade propagate through shared dependencies?"

[ASK] Check-my-math question: "Is that a fair picture of where this stands? What am I missing or getting wrong about your environment?" [PAUSE] after it. Silence is fine. Whatever they say, write it down visibly.

DERAIL-PREP: If someone says "we already have a plan / Nx / internal tooling for this" — don't argue. "That's great — Devin doesn't replace that plan, it executes it faster. Let me show you what I mean in the demo."

[RESEARCH BEFORE MEETING]
- Confirm public info on BofA's front-end stack (job postings are a goldmine: BofA postings frequently list Angular for digital banking roles — pull 2–3 current examples so 'you're an Angular shop' is grounded, not guessed).
- Check whether BofA has publicly discussed extended LTS vendors (HeroDevs NES etc.) — if they've bought extended support, reframe urgency around cost of extension rather than policy violation.""")

# ---------------- Slide 4: Why BofA is a strong candidate ----------------
s = add_slide()
header(s, "Why this conversation, here", "Why Bank of America is a strong candidate", "4")
c_w, c_h = 3.95, 2.6
card(s, 0.55, 1.8, c_w, c_h, "The scale", [
    "~$13B annual tech spend",
    "One of the largest engineering orgs anywhere",
], body_size=14)
card(s, 4.7, 1.8, c_w, c_h, "The AI posture", [
    "Erica · AI patent leader",
    "Exec buy-in for AI already exists",
], body_size=14)
card(s, 8.85, 1.8, c_w, c_h, "The environment", [
    "Regulated, audit-first culture",
    "Devin's PR + log model fits it natively",
], body_size=14)
rect(s, 0.55, 4.75, 12.25, 0.95, LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1)
tb(s, 0.85, 4.98, 11.7, 0.65, [[("The bigger the estate, the bigger the payoff from doing this systematically.", {"bold": True, "color": NAVY})]], size=16)
footer(s)
notes(s, """This slide answers "why are we the ones hearing this pitch?" — and flatters honestly, with facts.

Talk track: "Three reasons this conversation makes more sense here than almost anywhere else. First, scale: at ~$13B a year in tech, even a small percentage of engineering capacity reclaimed from maintenance is enormous in absolute terms. Second, you're not AI-skeptics — Erica, the patent portfolio, your leadership's public statements — the organizational muscle to adopt this already exists. Third, and counterintuitively: your regulated environment is an advantage. Devin works through pull requests and leaves a complete session log — it was built for places where every change needs a paper trail."

[ASK] "Where do you feel the maintenance-vs-greenfield squeeze most acutely right now?" — their answer tells you which use cases to propose in next steps. Write it down.

DERAIL-PREP: If someone says "scale cuts both ways — our estate is too messy for this": "That's exactly why we start with a bounded pilot on one slice, not a big-bang rollout. Messy is normal; we'll see how Devin handles it on your code."

[RESEARCH BEFORE MEETING]
- Verify the tech-spend figure and Erica stats from the latest earnings call (same items as the agenda-slide checklist).
- Find 1–2 public BofA exec quotes on GenAI for developers to cite verbatim.""")

# ---------------- Slide 5: What is Devin ----------------
s = add_slide()
header(s, "The 2-minute version", "What is Devin?", "5")
tb(s, 0.55, 1.5, 12.25, 0.6, [[("An AI software engineer", {"bold": True, "color": NAVY, "size": 18}), (" — you assign the task; it plans, codes, tests, and ships a PR.", {"size": 18})]])
c_w = 3.95
card(s, 0.55, 2.35, c_w, 2.45, "Works like an engineer", [
    "Own shell, editor, browser",
    "Runs builds & tests as it goes",
], body_size=14)
card(s, 4.7, 2.35, c_w, 2.45, "Fits your SDLC", [
    "Everything lands as a PR",
    "Your CI, your approvals — unchanged",
], body_size=14)
card(s, 8.85, 2.35, c_w, 2.45, "Built for this work", [
    "Migrations, upgrades, refactors",
    "Parallel sessions · full audit log",
], body_size=14)
rect(s, 0.55, 5.05, 12.25, 0.95, NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1)
tb(s, 0.85, 5.28, 11.7, 0.65, [[("Your engineers stop being typists and become reviewers.", {"bold": True, "color": WHITE})]], size=16)
footer(s)
notes(s, """Keep this SHORT — five minutes max. Execs don't want an architecture lecture; they want to know what it is, where it fits, and what stays under their control.

Key distinction to land: "You've probably all used or evaluated Copilot-style tools. Those make an engineer type faster. Devin is different — you hand it the whole task. 'Upgrade this repo from Angular 14 to 18, keep tests green' — and it comes back with reviewed-ready PRs."

Compliance hook for the Security Engineer: "Everything Devin does is captured in a session log — every command, every file change, every test run. For a bank, that's a better audit trail than most human engineers leave behind."

[ASK] calibration question: "What's your mental model of tools like this — copilot, contractor, or something else? I'll calibrate to that." It surfaces skepticism early, on your terms, instead of mid-demo.

DERAIL-PREP:
- "We tried AI coding tools and they hallucinate": "Fair — most tools generate code and hope. Devin runs the code. It compiles, executes tests, and iterates until they pass. You'll see that live in about 60 seconds."
- "What model is it? Where does our code go?": Don't improvise. "Great question — deployment options include VPC and data-isolation controls; I'll bring our security team for a proper review as part of next steps." [RESEARCH: know Cognition's current enterprise deployment/data-residency story cold before the meeting.]""")

# ---------------- Slide 6: Stakeholder concerns ----------------
s = add_slide()
header(s, "Addressing the elephants in the room", "What each of you is probably thinking", "6")
c_w, c_h = 3.95, 3.6
card(s, 0.55, 1.62, c_w, c_h, "VP of Engineering", [
    ("\u201CIs this a time-saver or a babysitting job?\u201D", {"italic": True, "color": GRAY}),
    [("→ ", {"bold": True, "color": RED}), ("Reviewing ≪ writing", {})],
    [("→ ", {"bold": True, "color": RED}), ("Roadmap keeps moving", {})],
], body_size=13.5)
card(s, 4.7, 1.62, c_w, c_h, "Security Engineer", [
    ("\u201CAn autonomous agent in a bank's codebase?\u201D", {"italic": True, "color": GRAY}),
    [("→ ", {"bold": True, "color": RED}), ("Nothing merges without human approval", {})],
    [("→ ", {"bold": True, "color": RED}), ("Every action logged & replayable", {})],
], body_size=13.5)
card(s, 8.85, 1.62, c_w, c_h, "Chief Architect", [
    ("\u201CWill it respect our standards?\u201D", {"italic": True, "color": GRAY}),
    [("→ ", {"bold": True, "color": RED}), ("Your conventions, applied consistently", {})],
    [("→ ", {"bold": True, "color": RED}), ("Library validated before downstream", {})],
], body_size=13.5)
footer(s)
notes(s, """This slide is your engagement insurance policy. Naming each person's likely concern out loud — before they raise it — earns trust and takes the wind out of ambushes.

Delivery tip: address each column TO that person, by name, with eye contact. "For you, [VP name], my guess is the question is team capacity..." Then [ASK]: "I took a guess at what each of you cares about. Did I get it roughly right — and what did I miss? I'd rather demo to your real concerns than my imagined ones." Let each of the three respond. Budget 3–4 minutes; this is where buy-in starts.

DERAIL-PREP (opinionated execs):
- If one person dominates: "That's a thread I want to pull on — but I want to make sure I hear [other person]'s take before the demo, since it changes what I show."
- If someone is openly skeptical/hostile: don't defend, redirect to evidence. "Totally fair to be skeptical. Can I earn the benefit of the doubt with the demo? If it doesn't hold up, tear it apart."
- If they push into pricing/procurement early: "I'll cover commercial structure in next steps — short version: we typically start with a bounded pilot so you're not buying on faith."

[RESEARCH BEFORE MEETING]
- Get the attendees' actual names/backgrounds (LinkedIn) — tailor each column's phrasing to how their org actually talks.
- Check BofA's public statements on responsible-AI governance (they have an internal AI governance framework) — mirror their language when talking about human oversight.""")

# ---------------- Slide 7: Demo ----------------
s = add_slide()
header(s, "Enough slides — let's watch it work", "Live demo: Devin migrates a BofA-shaped banking app, 14 → 18", "7")
card(s, 0.55, 1.62, 5.3, 4.3, "A codebase shaped like yours", [
    [("Shared component library", {"bold": True})],
    [("SSO/MFA · analytics SDK · financial data", {})],
    [("Custom design system on Angular Material", {})],
    ("Mock code, real architecture.", {"italic": True, "color": GRAY}),
], body_size=14)
card(s, 6.05, 1.62, 6.75, 4.3, "What Devin will do", [
    [("1. ", {"bold": True, "color": RED}), ("Baseline — build + tests green", {})],
    [("2. ", {"bold": True, "color": RED}), ("Stepwise ng update 14→15→16→17→18", {})],
    [("3. ", {"bold": True, "color": RED}), ("Modernize 5 legacy patterns", {})],
    [("4. ", {"bold": True, "color": RED}), ("Validate — zero legacy patterns remain", {})],
    [("5. ", {"bold": True, "color": RED}), ("Deliver — reviewable PRs", {})],
], body_size=14)
footer(s)
notes(s, """Switch to the live demo here (the !bofa_demo playbook flow). This slide is the map you leave up briefly so they know what they're about to see, then you go live.

While it runs, invite interaction: "Shout out anything you'd want to inspect — a diff, a test, a log. This is your demo, not my script."

Narrate value, not syntax: each of the 5 patterns has a business hook —
1. HttpParams → "API evolution handled systematically, not by grep-and-pray."
2. Standalone components → "Adopting Angular 18 the right way, not just compiling on it."
3. TestBed → "Your test suite migrates too — most manual upgrades leave tests rotting."
4. any types in financial models → point at the Security Engineer: "Type safety on money-handling code. This is a defect class disappearing."
5. *ngFor → @for → point at the Architect: "Template modernization applied identically across the whole library — consistency at scale."

Pattern 4 and the downstream-teams story are your strongest moments — slow down there.

[PAUSE] moments: after the baseline goes green ("that's the safety net"), and after the final validation scan ("zero legacy patterns — that's the compliance evidence your audit file wants").

DERAIL-PREP:
- Demo breaks: stay calm, it's rehearsed — you have git-reset recovery branches and a recorded run as backup. "Live demos, right? Here's the recorded run while I reset."
- "This demo repo is toy-sized; ours is 100x bigger": "Completely true — and that's the point of the pilot: we prove it on a real slice of YOUR codebase in week one. Scale is Devin's advantage: it parallelizes; humans don't."
- "Devin wrote this demo repo, of course it can fix it": "Fair — which is why next step is your code, not mine." """)

# ---------------- Slide 8: How Devin solves it ----------------
s = add_slide()
header(s, "From demo to your reality", "How this maps to BofA's actual migration", "8")
steps = [
    ("Week 0", "Scope & safety rails", "Security review · success criteria"),
    ("Weeks 1–2", "Shared library first", "Root of the dependency tree — your team reviews every PR"),
    ("Weeks 2–3", "Downstream in parallel", "One Devin per consuming app — the coordination tax collapses"),
    ("Ongoing", "Evidence as you go", "Logs + tests + scans = audit file, automatically"),
]
y = 1.62
for i, (when, t, d) in enumerate(steps):
    rect(s, 0.55, y, 12.25, 1.02, LIGHT if i % 2 == 0 else WHITE, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1)
    rect(s, 0.75, y + 0.18, 1.5, 0.66, NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.18)
    tb(s, 0.75, y + 0.33, 1.5, 0.4, when, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    tb(s, 2.5, y + 0.12, 10.0, 0.4, t, size=14.5, color=NAVY, bold=True)
    tb(s, 2.5, y + 0.5, 10.0, 0.45, d, size=12, color=DARK)
    y += 1.12
rect(s, 0.55, 6.2, 12.25, 0.72, NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1)
tb(s, 0.85, 6.33, 11.7, 0.5, [[("Devin changes who does the labor — never who's in control.", {"bold": True, "color": WHITE})]], size=15)
footer(s)
notes(s, """Bridge from demo to their world: "What you just saw on five patterns is exactly what runs across your real repos — the difference is scale, and scale is where this gets interesting."

The killer point for the Architect is step 3: parallelization. "A human team upgrades one repo at a time and coordinates by meeting. Devin runs one session per downstream app simultaneously. Your critical path goes from 'sum of 12 teams' calendars' to 'longest single review cycle.'"

For the Security Engineer, step 4: "You don't build the audit story at the end — every session IS the audit story."

[ASK] the Architect directly: "If we upgraded the shared library first and handed each consuming team a passing PR, what would still worry you?" — whatever they say becomes a pilot success criterion. Write it down.

DERAIL-PREP:
- "Our repos are on internal GitHub Enterprise / restricted network": "Supported — deployment specifics are exactly what the week-0 security review covers." [RESEARCH: confirm Cognition's current GHE/VPC/network-isolation options before the meeting so you don't overpromise.]
- "Who's liable if Devin introduces a bug?": "Same person as today — the reviewing engineer and your existing release process. Devin doesn't get commit rights your policies wouldn't give a contractor." """)

# ---------------- Slide 9: End result ----------------
s = add_slide()
header(s, "What good looks like", "The end state: BofA running on Devin", "9")
c_w, c_h = 3.95, 2.5
card(s, 0.55, 1.62, c_w, c_h, "This migration", [
    "Angular 18 before the deadline",
    "12+ teams unbroken",
    "Audit trail, automatic",
], body_size=14)
card(s, 4.7, 1.62, c_w, c_h, "Your engineers", [
    "Back on revenue features",
    "Reviewing, not hand-migrating",
], body_size=14)
card(s, 8.85, 1.62, c_w, c_h, "The bigger unlock", [
    "Angular 19, 20… become routine",
    "Same machinery: cloud migration, test coverage",
], body_size=14)
rect(s, 0.55, 4.4, 12.25, 1.15, LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1)
tb(s, 0.85, 4.62, 11.7, 0.85, [
    [("\u201CWe met a hard compliance deadline without pausing the roadmap.\u201D", {"italic": True, "bold": True, "color": NAVY})],
], size=17)
footer(s)
notes(s, """Paint the after-picture in THEIR terms, one beat per persona:
- VP Eng: capacity back + deadline met without roadmap carnage.
- Security: policy violation resolved, plus a standing audit trail for every future change.
- Architect: consistency enforced mechanically across the estate; migrations become boring (the highest compliment in architecture).

The third card is deliberately the expansion story — mention the other two use cases from their own list (cloud migration, test coverage) briefly: "And to be direct, the same machinery applies to the Lambda migration and the OCC test-coverage push you mentioned. But let's earn that with this one first." Referencing their own pre-shared initiatives shows you listened.

[ASK] gut-check: "If this end state showed up on time, what would it be worth to your org? I'm asking seriously — it should anchor how we scope the pilot." [PAUSE] and let them put a number or a feeling on it. Whatever value they articulate becomes YOUR value framing in follow-up emails — their words, not yours.""")

# ---------------- Slide 10: Next steps ----------------
s = add_slide()
header(s, "Leaving with a decision, not a deck", "Proposed next steps", "10")
steps = [
    ("1", "This week", "Pick 3 use cases together", "Start with work like today's demo — upgrades, migrations, test debt"),
    ("2", "Next week", "Security & data review", "Your security team + ours · define success criteria"),
    ("3", "Weeks 2–4", "Run a pilot for you", "Devin works one real use case; your engineers review; we measure"),
    ("4", "Week 4", "Go / no-go readout", "Results in front of this room — then scale to the other two"),
]
y = 1.62
for n, when, t, d in steps:
    rect(s, 0.55, y, 12.25, 1.02, LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1)
    rect(s, 0.72, y + 0.22, 0.58, 0.58, RED, shape=MSO_SHAPE.OVAL)
    tb(s, 0.72, y + 0.3, 0.58, 0.44, n, size=17, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    tb(s, 1.55, y + 0.12, 1.55, 0.75, when, size=11.5, color=RED, bold=True)
    tb(s, 3.1, y + 0.12, 9.4, 0.4, t, size=14.5, color=NAVY, bold=True)
    tb(s, 3.1, y + 0.5, 9.4, 0.45, d, size=11.5, color=DARK)
    y += 1.12
tb(s, 0.55, 6.25, 12.25, 0.6, [[("The ask:  ", {"bold": True, "color": RED, "size": 14}), ("three use cases, one pilot, four weeks, your success criteria.", {"size": 14, "bold": True, "color": NAVY})]])
footer(s)
notes(s, """End with a specific, low-risk ask and assign owners in the room — that's what turns a nice meeting into a pipeline.

Talk track: "Here's what I'd propose. First, we sit down together and pick three use cases — the natural starting points are work like what you just watched: framework upgrades, migrations, test-coverage debt — but you know your backlog; whatever you'd never staff greenfield engineers on is a candidate. Then a security and data review so your team is comfortable. Then we run a pilot on ONE of the three — a real repo, four weeks, measured against criteria you define. If it earns its keep, we scale to the other two. If it doesn't, you've spent almost nothing finding out."

Why 3 use cases but 1 pilot: picking three gets them mentally invested in a roadmap (not a one-off experiment), while piloting one keeps the ask small enough to say yes to today.

Assign ownership live: use-case shortlist → VP Eng (it's their backlog); security review → Security Engineer; pilot repo + success criteria → Chief Architect. People who own a step show up to the next meeting.

[ASK] the closing question and get a name and a date before you leave the room. "Who should I follow up with, and does end of this week work for the security call?"

DERAIL-PREP:
- "We need to run this through procurement/vendor risk first": "Understood — that's exactly why step 1 is the security review; it feeds your vendor-risk process. Can we start that in parallel?"
- "Send us materials and we'll get back to you": counter gently with a smaller yes: "Happy to — can I also pencil the security call so materials land with momentum? Easy to cancel."
- If the room is split (one enthusiast, one skeptic): scope the pilot around the skeptic's concern explicitly — "Let's make [skeptic]'s question the primary success criterion."

Post-meeting: send notes + promised materials within 24–48h, restate the owners and dates they agreed to.""")

# ---------------- Slide 11: Appendix divider ----------------
s = add_slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, 3.9, SW, 0.06, RED)
tb(s, 0.9, 2.9, 11.5, 0.9, "Appendix — Prep Material (not presented)", size=34, color=WHITE, bold=True)
tb(s, 0.9, 4.15, 11.0, 1.0, [
    ("A. Objection handling — if the room pushes back", {"size": 16, "color": RGBColor(0xC9,0xD4,0xEC)}),
    ("B. Pre-meeting research checklist", {"size": 16, "color": RGBColor(0xC9,0xD4,0xEC)}),
])
notes(s, "These slides are for your prep only — hide them (right-click → Hide slide) or delete before presenting. Keep them in your back pocket; if the meeting goes deep on an objection, you'll have rehearsed the answer.")

# ---------------- Slide 12: Objection handling ----------------
s = add_slide()
header(s, "Appendix A", "If the room pushes back: objections & responses", "A")
objs = [
    ("\u201CAI writing code for a bank? Our regulators will never allow it.\u201D",
     "Devin changes who types, not who approves. Every change goes through your existing review, CI, and release gates — the same controls you'd apply to a contractor. And it leaves a better audit trail than a human."),
    ("\u201CWe already have Copilot / internal AI tools.\u201D",
     "Complementary, not competing. Copilot accelerates an engineer who's already assigned. Devin absorbs whole task categories nobody wants assigned — migrations, upgrades, test debt."),
    ("\u201COur codebase is far bigger and messier than any demo.\u201D",
     "Agreed — that's why the pilot is on your code, your criteria, in four weeks. Scale actually favors Devin: it parallelizes across repos; humans don't."),
    ("\u201CWhat about hallucinations / wrong code?\u201D",
     "Devin executes what it writes — builds, tests, and iterates until green before opening a PR. Wrong code fails your test suite and never reaches a human's desk, let alone production."),
    ("\u201CIs this going to replace my engineers?\u201D",
     "It replaces the work your engineers resent, not the engineers. The pitch to your team: nobody hand-migrates *ngFor loops anymore; everyone reviews, designs, and ships features."),
    ("\u201CWhere does our code go? Data security?\u201D",
     "Answered properly in the week-0 security review — enterprise deployment and data-isolation options exist. Don't accept a hand-wave from me; bring your security team."),
]
y = 1.55
for q, a in objs:
    rect(s, 0.55, y, 12.25, 0.82, LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.1)
    tb(s, 0.8, y + 0.06, 4.55, 0.72, [(q, {"italic": True, "bold": True, "color": NAVY})], size=10.5, line_spacing=0.95)
    tb(s, 5.5, y + 0.06, 7.15, 0.72, a, size=10, color=DARK, line_spacing=0.95)
    y += 0.9
footer(s)
notes(s, """General derailment playbook (for opinionated execs):
1. Never argue — agree with the valid kernel, then redirect to evidence ("fair concern; here's what we can measure").
2. Park with a return time: "Great question — that's next-steps territory; give me 8 minutes and I'll answer it properly." Then ACTUALLY return to it, by name.
3. If two execs start debating each other, don't referee — extract the requirement: "Sounds like the pilot needs to prove X to satisfy both of you. Noted."
4. If someone tries to turn it into a vendor bake-off ("how do you compare to X?"): "Happy to do a feature comparison async — in the room, the fastest way to compare is a pilot on the same repo."
5. Protect the demo time. If discovery is running long, cut slides 8–9 short, never the demo — the demo is what they'll remember.
6. If the meeting is fully derailed by one topic (e.g., security), embrace it: "Let's make this the security deep-dive, and I'll get you the demo recording + a follow-up demo slot." A derailed meeting where one stakeholder gets deeply satisfied still advances the deal.""")

# ---------------- Slide 13: Research checklist ----------------
s = add_slide()
header(s, "Appendix B", "Pre-meeting research checklist (do before the room)", "B")
items_l = [
    "BofA tech budget: verify latest figure (~$12–13B total, ~$3.5–4B new initiatives) from most recent earnings call.",
    "Erica + AI posture: current interaction stats, AI patent counts, exec quotes (Aditya Bhasin, Hari Gopalkrishnan) on GenAI for developers.",
    "Angular evidence: pull 2–3 live BofA job postings listing Angular for digital banking — grounds the 'your stack' claim.",
    "Angular 14 EOL specifics: community EOL Nov 2023; check if BofA uses extended LTS (HeroDevs NES etc.) — changes urgency framing.",
    "Attendees: names, LinkedIn backgrounds, org context for all three; adjust the personas slide phrasing to match.",
]
items_r = [
    "Cognition enterprise story: deployment options (VPC/isolation), data residency, SOC 2 / compliance posture — know it cold for security Qs.",
    "Public financial-services references or case studies you're allowed to cite (banks, regulated industries).",
    "BofA AI governance: their published responsible-AI framework language — mirror it when discussing oversight.",
    "Pricing/pilot structure: know the current pilot commercial model so 'what does it cost' gets a confident one-liner.",
    "Rehearse demo failure recovery: git reset branches + recorded backup run tested on the presentation machine.",
]
card(s, 0.55, 1.62, 6.0, 4.6, "About Bank of America", items_l, body_size=11.5)
card(s, 6.8, 1.62, 6.0, 4.6, "About Cognition / logistics", items_r, body_size=11.5)
tb(s, 0.55, 6.35, 12.25, 0.5, [[("Rule of thumb: ", {"bold": True, "color": RED, "size": 12}), ("every number you say out loud in that room should have a source you checked within the last month.", {"size": 12, "italic": True})]])
footer(s)
notes(s, """This slide is your homework list — everything the deck references that needs verification or that I (Devin) couldn't confirm as current fact. Delete before presenting.

Highest priority items:
1. The ~$13B tech budget and Erica stats — these are your two 'I did my homework' credibility numbers; get them current and sourced.
2. Cognition's enterprise security/deployment answers — the Security Engineer WILL ask, and a fumbled answer undoes the whole demo.
3. Attendee names/backgrounds — the personas slide is dramatically stronger delivered to real people by name.
4. The extended-LTS question — if BofA already bought Angular extended support, your urgency pitch shifts from 'policy violation' to 'you're paying rent on dead software; here's the exit.'""")

out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "BofA_x_Cognition_Devin_Demo.pptx")
prs.save(out)
print("slides:", len(prs.slides._sldIdLst))
